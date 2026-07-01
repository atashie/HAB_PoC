#!/usr/bin/env python3
"""Build HAB_landscape.pptx from the traceable data JSON + generated figures (python-pptx).

Slides: (1) Tools lead, (2) Tools appendix full table, (3) Models lead, (4) Models appendix
full table, (5) Methods & how-to-read (quality-tier rubric + provenance). Every table row
carries source citation keys; figures are the generated PNGs; speaker notes on each slide.
16:9. Run build_figures.py first. Deterministic.
"""
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
DATA = os.path.join(HERE, "data")
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "HAB_landscape.pptx")

tools = json.load(open(os.path.join(DATA, "tools.json"), encoding="utf-8"))["tools"]
models = json.load(open(os.path.join(DATA, "models.json"), encoding="utf-8"))["models"]

# palette
NAVY = RGBColor(0x27, 0x4C, 0x77)
BLUE = RGBColor(0x3B, 0x6E, 0xA5)
ORANGE = RGBColor(0xE0, 0x7B, 0x39)
GREEN = RGBColor(0x4C, 0x95, 0x6C)
LGREY = RGBColor(0xEE, 0xF1, 0xF5)
MGREY = RGBColor(0xD6, 0xDC, 0xE3)
DGREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)

TIER_COLOR = {"established-operational": GREEN, "operational-unbenchmarked": BLUE,
              "vendor-self-report": ORANGE, "research-grade": NAVY}
TIER_SHORT = {"established-operational": "Established / operational", "operational-unbenchmarked": "Operational (unbenchmarked)",
              "vendor-self-report": "Vendor self-report", "research-grade": "Research-grade"}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def set_run(r, text, size, bold=False, color=BLACK, italic=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"


def title_bar(slide, kicker, title, accent=NAVY):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    _, tf = textbox(slide, Inches(0.45), Inches(0.10), SW - Inches(0.9), Inches(0.95), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p.add_run(), kicker.upper(), 11, bold=True, color=RGBColor(0xCF, 0xDD, 0xEC))
    p2 = tf.add_paragraph()
    set_run(p2.add_run(), title, 24, bold=True, color=WHITE)
    return bar


def footnote(slide, text):
    _, tf = textbox(slide, Inches(0.45), SH - Inches(0.42), SW - Inches(0.9), Inches(0.36))
    set_run(tf.paragraphs[0].add_run(), text, 8, italic=True, color=DGREY)


def add_image_fit(slide, path, l, t, w, h):
    """Add image scaled to fit within (w,h) box, centered, preserving aspect ratio."""
    from PIL import Image
    iw, ih = Image.open(path).size
    box_ar = w / h
    img_ar = iw / ih
    if img_ar > box_ar:
        nw = w; nh = int(w / img_ar)
    else:
        nh = h; nw = int(h * img_ar)
    ll = l + (w - nw) // 2
    tt = t + (h - nh) // 2
    return slide.shapes.add_picture(path, ll, tt, nw, nh)


def style_cell(cell, text, size=9, bold=False, color=BLACK, fill=None, align=PP_ALIGN.LEFT):
    cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size, bold=bold, color=color)


def build_table(slide, l, t, w, headers, rows, col_w, header_fill=NAVY, body_size=8.5, header_size=9, row_h=Inches(0.3)):
    nrow = len(rows) + 1
    ncol = len(headers)
    tbl_shape = slide.shapes.add_table(nrow, ncol, l, t, w, row_h * nrow)
    tbl = tbl_shape.table
    tbl.first_row = False; tbl.horz_banding = False
    total = sum(col_w)
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Emu(int(w * cw / total))
    for j, htext in enumerate(headers):
        style_cell(tbl.cell(0, j), htext, size=header_size, bold=True, color=WHITE, fill=header_fill,
                   align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows, start=1):
        band = LGREY if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            fill = band
            color = BLACK
            if isinstance(val, tuple):  # (text, fillcolor, textcolor)
                text, fill, color = val
            else:
                text = val
            style_cell(tbl.cell(i, j), text, size=body_size, color=color, fill=fill)
    return tbl_shape


def norm_access(a):
    a = a.lower()
    return "Freemium" if "freemium" in a else "Paid" if "paid" in a else "Free" if "free" in a else a


def keys_str(keys, n=3):
    return ", ".join(keys[:n]) + ("…" if len(keys) > n else "")


def trunc(s, n):
    s = s or ""
    if len(s) <= n:
        return s
    cut = s[:n].rsplit(" ", 1)[0]
    return cut + "…"


import re as _re


def compact_horizon(t):
    """Short horizon label for dense tables (full phrasing lives in the JSON + notes)."""
    h = t["horizon"].lower()
    is_fc = "forecast" in h
    is_now = "nowcast" in h
    lead = ""
    m = _re.search(r"(\d+\s*(?:-|\s)?(?:day|hr|hour|week))", h)
    if m:
        lead = " " + m.group(1).replace(" ", "").replace("-", "‑")
    if is_fc and is_now:
        return "Nowcast + forecast" + lead
    if is_fc:
        return "Forecast" + lead
    return "Nowcast"


SIGNAL_SHORT = {"fusion": "Fusion", "satellite": "Satellite", "hybrid": "Hybrid", "n/a": "Review",
                "in-situ": "In-situ", "in-situ + weather": "In-situ+wx", "in-situ + physical": "In-situ+phys",
                "in-situ + loads": "In-situ+loads"}


# ============================== SLIDE 1: TOOLS LEAD ==============================
s = add_slide()
title_bar(s, "Slide 1 · What's available today", "HAB monitoring & forecasting tools you can use now", NAVY)
add_image_fit(s, os.path.join(FIG, "fig1_tools_panel.png"), Inches(0.3), Inches(1.22), Inches(4.4), Inches(5.6))

# takeaways (right, expanded)
_, tf = textbox(s, Inches(4.95), Inches(1.35), Inches(8.15), Inches(3.4))
set_run(tf.paragraphs[0].add_run(), "What to take away", 14, bold=True, color=NAVY)
takeaways = [
    "12 freshwater tools surveyed — 9 are free & public; all are US federal + state programs + a few private services. No NGO or international operational freshwater tools exist.",
    "Only 4 of the 12 forecast ahead (red outline) — every forecaster is federal (EPA, NOAA) or private; all state / local tools are nowcast / observation only.",
    "The forecasters split by cost: EPA's 7-day forecast is free; the private forecasters (BlueGreen, CyanoLakes) are paid and publish no validation.",
    "Satellite (EPA CyAN) is the only near-national signal — but it sees biomass, not toxin, and can't resolve small lakes or shoreline scums.",
    "Methods span satellite indices, in-situ monitoring, a physics model and ML. 'Data fusion' (satellite + in-situ + weather, used by EPA, NOAA & BlueGreen) is a data strategy run by one of these methods — not a method of its own.",
]
for tk in takeaways:
    p = tf.add_paragraph()
    set_run(p.add_run(), "•  " + tk, 11, color=BLACK)
    p.space_after = Pt(6)

# forecast-only tools table (right, below takeaways)
_, tf2 = textbox(s, Inches(4.95), Inches(4.82), Inches(8.15), Inches(0.32))
set_run(tf2.paragraphs[0].add_run(), "The 4 forward-looking (forecast) tools — details", 11, bold=True, color=NAVY)
fc = [t for t in tools if "forecast" in t["horizon"].lower()]
NAME_S = {"epa7day": "EPA 7-day forecast", "noaa_erie": "NOAA Lake Erie", "bluegreen": "BlueGreen BGi", "cyanolakes": "CyanoLakes"}
METH_S = {"empirical-RS": "Satellite", "statistical-ML": "Statistical/ML", "mechanistic": "Physics", "in-situ-sampling": "In-situ", "fusion": "Fusion"}
REFRESH_S = {"epa7day": "Weekly", "noaa_erie": "2x/wk (season)", "bluegreen": "n/r", "cyanolakes": "n/r"}
rows = []
for t in fc:
    cov = t["coverage"].replace(", CONUS", "").replace(" (commercial)", "")
    res = t["spatial_res_m"].replace(" (satellite input)", "")
    rows.append([NAME_S.get(t["id"], t["name"]), METH_S.get(t["model_type"], t["model_type"]),
                 REFRESH_S.get(t["id"], t["refresh"]), cov, res])
build_table(s, Inches(4.95), Inches(5.18), Inches(8.15), ["Tool", "Method", "Refresh", "Coverage", "Spatial res."], rows,
            col_w=[2.2, 1.5, 1.7, 2.6, 1.4], body_size=9.5, header_size=9.5, row_h=Inches(0.4))
footnote(s, "Freshwater cyanobacteria only. Every tool traces to a cited source (appendix + Research/REFERENCES.md). Raw counts, small N — orientation, not statistics.")
s.notes_slide.notes_text_frame.text = (
    "Slide 1 orients the panel to what exists TODAY. The three barplots: who makes them (state+federal dominate; no NGO/intl), "
    "cost (9 of 12 free), and method. The bottom-right table isolates the only 4 forward-looking (forecast) tools — EPA's free "
    "7-day product plus three paid/vendor services (two unvalidated). Most other tools only nowcast. Satellite CyAN is the only "
    "near-national signal but detects biomass (not toxin) and can't resolve small lakes — motivating a fused approach. No "
    "NGO/international operational freshwater tool surfaced in our 261-source review.")

# ============================== SLIDE 2: TOOLS APPENDIX ==============================
s = add_slide()
title_bar(s, "Slide 1 · Appendix", "Full tool inventory — all dimensions", BLUE)
headers = ["Tool", "Org", "Method", "Horizon", "Refresh", "Spatial res.", "Access", "Year", "Quality tier", "Src"]
col_w = [2.5, 1.0, 1.15, 1.7, 1.7, 2.1, 0.95, 0.85, 1.85, 1.05]
rows = []
for t in tools:
    tier = t["quality_tier"]
    rows.append([
        t["name"].replace(" + bulletin", "").replace("EPA experimental 7-day cyanoHAB forecast", "EPA 7-day forecast"),
        {"federal": "Federal", "state-local": "State/local", "private": "Private"}.get(t["org_type"], t["org_type"]),
        {"empirical-RS": "Satellite", "statistical-ML": "Statistical/ML", "mechanistic": "Physics", "in-situ-sampling": "In-situ", "fusion": "Fusion"}.get(t["model_type"], t["model_type"]),
        compact_horizon(t),
        trunc(t["refresh"], 32),
        trunc(t["spatial_res"], 40),
        (norm_access(t["access"]), ORANGE if "paid" in t["access"].lower() else (GREEN if norm_access(t["access"]) == "Free" else BLUE), WHITE),
        t["year"],
        (TIER_SHORT[tier], TIER_COLOR[tier], WHITE),
        keys_str(t["keys"], 2),
    ])
build_table(s, Inches(0.25), Inches(1.30), Inches(12.85), headers, rows, col_w, body_size=8, header_size=8.5, row_h=Inches(0.44))
footnote(s, "Quality tiers are evidence-anchored — the cited evidence basis for each is in presentation/data/tools.json & the source dossiers; rubric on final slide. 'n/r' = not reported. Src → Research/REFERENCES.md.")
notes = ["Full tool reference. Evidence basis behind each quality tier (read if challenged):"]
for t in tools:
    notes.append(f"- {t['name']}: [{TIER_SHORT[t['quality_tier']]}] {t['quality_basis']}")
s.notes_slide.notes_text_frame.text = "\n".join(notes)

# ============================== SLIDE 3: MODELS LEAD ==============================
s = add_slide()
title_bar(s, "Slide 2 · The research frontier", "Cutting-edge HAB models — and what actually drives them", NAVY)
add_image_fit(s, os.path.join(FIG, "fig2_feature_frequency.png"), Inches(0.25), Inches(1.25), Inches(8.25), Inches(5.15))

_, tf = textbox(s, Inches(8.6), Inches(1.30), Inches(4.55), Inches(3.9))
set_run(tf.paragraphs[0].add_run(), "What to take away", 13, bold=True, color=NAVY)
tks = [
    "Water temperature & nutrients are the near-universal inputs; water temperature is most often the #1 predictor — but that's a correlational, model-internal signal (often a seasonal proxy), not a proven cause.",
    "Exceptions matter: some models rank phycocyanin or nutrients first, and one fusion model's top 'feature' was geolocation — a non-transferable artifact.",
    "Most of the 14 frontier models are simple & explainable (tree/boosting) and in-situ-driven; deep learning leads only in aggregate reviews.",
    "Persistence is a stiff baseline (week-to-week autocorrelation ≈0.90) — a real tool must beat it.",
]
for tk in tks:
    p = tf.add_paragraph()
    set_run(p.add_run(), "•  " + tk, 10, color=BLACK)
    p.space_after = Pt(5)

_, tf2 = textbox(s, Inches(8.6), Inches(5.5), Inches(4.55), Inches(0.3))
set_run(tf2.paragraphs[0].add_run(), "Flagship models & their #1 predictor", 10, bold=True, color=DGREY)
flag = [m for m in models if m["id"] in ("M1", "M3", "M6", "M9", "M11")]
rows = []
for m in flag:
    rows.append([trunc(m["name"].split(",")[0], 26), (m.get("top_features") or ["—"])[0].replace("_", " ")[:16]])
build_table(s, Inches(8.6), Inches(5.8), Inches(4.55), ["Model", "Top predictor"], rows,
            col_w=[3.3, 1.7], body_size=8, row_h=Inches(0.2), header_fill=NAVY)
footnote(s, "Feature 'importance' = how often a feature is a model's top-ranked predictor (correlational, not causal). 10 of the 14 models report a feature list; the 168-study review (ACAD-068) is context, not counted as a model.")
s.notes_slide.notes_text_frame.text = (
    "Slide 2 answers 'what does the frontier look like and what drives it.' The hero chart: water temperature & nutrients "
    "are near-universal; water temperature is most often the #1 predictor, BUT stress this is correlational/model-internal "
    "(often a season proxy) and has real exceptions (phycocyanin in Müggelsee; geolocation in the Sentinel-2 fusion model — "
    "an explicitly non-transferable artifact). Bottom-left: frontier ML is mostly simple tree/boosting methods and in-situ-"
    "driven (Korea/China/Germany dominate); deep learning only leads in aggregate reviews. The persistence baseline (AR1≈0.90) "
    "is the number any PoC must beat — this directly frames our modeling choices.")

# ============================== SLIDE 4: MODELS APPENDIX ==============================
s = add_slide()
title_bar(s, "Slide 2 · Appendix", "Full model inventory (14 models + 1 review) — features emphasised", BLUE)
headers = ["Model", "Org", "Algorithm", "Forecasts", "Mode", "Signal", "Features used  (→ TOP predictor in bold text)", "Src"]
col_w = [2.5, 0.95, 2.15, 1.7, 0.95, 1.15, 4.7, 0.95]
rows = []
for m in models:
    _fl = lambda f: f.replace("_", " ").replace("phycocyanin signal", "phycocyanin (in-situ)")
    feats = ", ".join(_fl(f) for f in m["features"]) if m["features"] else "n/r"
    top = ", ".join(_fl(t) for t in m["top_features"]) if m["top_features"] else "—"
    feat_cell = feats + ("  → TOP: " + top if m["top_features"] else "")
    rows.append([
        trunc(m["name"].split(",")[0], 30),
        m["org_type"].replace(" / ", "/").replace("federal", "Fed").replace("academic", "Acad"),
        trunc(m["algorithm"], 32),
        trunc(m["target"], 22),
        m["mode"],
        SIGNAL_SHORT.get(m["data_signal"], m["data_signal"]),
        trunc(feat_cell, 96),
        keys_str(m["keys"], 2),
    ])
build_table(s, Inches(0.2), Inches(1.28), Inches(12.95), headers, rows, col_w, body_size=7, header_size=7.5, row_h=Inches(0.34))
add_image_fit(s, os.path.join(FIG, "fig3_models_context.png"), Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.4))
footnote(s, "The 168-study REVIEW (mode='review') is context, not counted among the 14 models. Feature importance is correlational / model-internal (not causal). 'n/r' = not reported. Skill metrics in each dossier & data/models.json. Src → Research/REFERENCES.md.")
notes = ["Full model table; the features column is the brief's emphasis. Reported skill per model (read if asked):"]
for m in models:
    notes.append(f"- {m['name']}: {m['skill']}")
notes.append("\nFlag M11's longitude top-feature as a non-transferable proxy, and M1's 49%-precision / sample-count caveat.")
s.notes_slide.notes_text_frame.text = "\n".join(notes)

# ============================== SLIDE 5: METHODS & RUBRIC ==============================
s = add_slide()
title_bar(s, "How to read these slides", "Method, quality rubric & provenance", DGREY)
_, tf = textbox(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.6))
set_run(tf.paragraphs[0].add_run(), "Where this comes from", 14, bold=True, color=NAVY)
for line in [
    "Built entirely from a 261-source public-literature review (peer-reviewed, federal, state, NGO/WHO, private).",
    "Every tool/model attribute traces to a citation key that resolves in Research/REFERENCES.md; unknowns are marked 'n/r', never guessed.",
    "Freshwater cyanobacteria scope; marine systems excluded (a few tools' marine parents are footnoted).",
    "Counts are raw (small N) — orientation, not statistics.",
]:
    p = tf.add_paragraph(); set_run(p.add_run(), "•  " + line, 11); p.space_after = Pt(7)

set_run(tf.add_paragraph().add_run(), "", 6)
p = tf.add_paragraph(); set_run(p.add_run(), "Feature 'importance' caveat", 13, bold=True, color=NAVY)
p = tf.add_paragraph()
set_run(p.add_run(), "‘Most important’ = most often a model's top-ranked predictor. This is a correlational, "
        "model-internal signal (often a seasonal proxy) — NOT a demonstrated cause. Rankings differ across models.", 11)

# rubric table on right
_, tf3 = textbox(s, Inches(6.85), Inches(1.35), Inches(6.0), Inches(0.4))
set_run(tf3.paragraphs[0].add_run(), "Evidence-anchored quality tiers", 14, bold=True, color=NAVY)
rubric = [
    [("Established / operational", GREEN, WHITE), "Public product with peer-reviewed AND out-of-sample validation + baseline (e.g. CyAN, EPA 7-day)"],
    [("Operational (unbenchmarked)", BLUE, WHITE), "Deployed/public but no skill metric stated in the sources (e.g. NOAA Erie, state dashboards)"],
    [("Vendor self-report", ORANGE, WHITE), "Commercial; validation claimed but no disclosed baseline/peer review (BlueGreen, CyanoLakes, EOMAP)"],
    [("Research-grade", NAVY, WHITE), "Peer-reviewed model with reported skill; not an operational public product"],
]
build_table(s, Inches(6.85), Inches(1.85), Inches(6.15), ["Tier", "What it means"], rubric,
            col_w=[2.1, 4.0], body_size=9.5, header_size=10, row_h=Inches(0.75))
_, tf4 = textbox(s, Inches(6.85), Inches(5.2), Inches(6.15), Inches(1.8))
set_run(tf4.paragraphs[0].add_run(), "Reproducible", 13, bold=True, color=NAVY)
p = tf4.add_paragraph()
set_run(p.add_run(), "Figures & tables regenerate from presentation/data/*.json via build_figures.py + build_deck.py. "
        "Data JSON carries the source keys for every cell.", 11)
footnote(s, "SePRO Lead Data Scientist case study · HAB landscape orientation · generated 2026-07-01 · freshwater cyanobacteria scope.")
s.notes_slide.notes_text_frame.text = ("Reference slide. Use if the panel asks 'how do you know' or 'what does cutting-edge mean here'. The tier rubric "
                                        "makes the subjective 'quality' call defensible; the feature-importance caveat pre-empts a correlation-vs-causation challenge.")

prs.save(OUT)
print("saved", OUT, "| slides:", len(prs.slides._sldIdLst))
